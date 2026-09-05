import os
import io
import mimetypes
import google.genai as genai
from google.api_core import exceptions as google_exceptions
from PIL import Image
from django.conf import settings
from .base import OCRProvider

# Module-level model caches
_CACHED_GEMINI = None
_CACHED_AZURE_CLIENT = None


def get_gemini_model():
    global _CACHED_GEMINI
    if _CACHED_GEMINI is None:
        api_key = getattr(settings, "GEMINI_API_KEY", os.getenv("GEMINI_API_KEY", ""))
        if api_key:
            genai.configure(api_key=api_key)
        _CACHED_GEMINI = genai.GenerativeModel('gemini-1.5-flash')
    return _CACHED_GEMINI


def get_azure_client():
    global _CACHED_AZURE_CLIENT
    if _CACHED_AZURE_CLIENT is None:
        endpoint = getattr(settings, "AZURE_DOCUMENT_INTELLIGENCE_ENDPOINT", os.getenv("AZURE_DOCUMENT_INTELLIGENCE_ENDPOINT", ""))
        key = getattr(settings, "AZURE_DOCUMENT_INTELLIGENCE_KEY", os.getenv("AZURE_DOCUMENT_INTELLIGENCE_KEY", ""))
        if endpoint and key:
            from azure.core.credentials import AzureKeyCredential
            from azure.ai.documentintelligence import DocumentIntelligenceClient
            _CACHED_AZURE_CLIENT = DocumentIntelligenceClient(
                endpoint=endpoint,
                credential=AzureKeyCredential(key)
            )
    return _CACHED_AZURE_CLIENT


class HybridOCRProvider(OCRProvider):
    def __init__(self):
        self.vision_prompt = (
            "Analyze this learning note/document image thoroughly:\n"
            "1. Transcribe all handwritten and printed text accurately.\n"
            "2. If diagrams, flowcharts, graphs, or formulas are present, clearly describe their structure, "
            "labels, data points, and relationships in structured text.\n"
            "3. Maintain logical reading order. Output clear, unified text without conversational filler."
        )

    def _resolve_file_path(self, file_or_path) -> str:
        """Extracts the absolute filesystem path from FieldFile or string."""
        if hasattr(file_or_path, 'path'):
            return file_or_path.path
        return str(file_or_path)

    # ==========================================================
    # AZURE AI DOCUMENT INTELLIGENCE FALLBACK
    # ==========================================================
    def _run_azure_fallback(self, image_bytes: bytes) -> str:
        """Runs Azure Document Intelligence prebuilt-layout model when Gemini quota is exceeded."""
        client = get_azure_client()
        if client is None:
            raise RuntimeError("[Error: Gemini API limit reached and Azure Document Intelligence is not configured.]")

        try:
            from azure.ai.documentintelligence.models import AnalyzeDocumentRequest
            poller = client.begin_analyze_document(
                model_id="prebuilt-layout",
                analyze_request=AnalyzeDocumentRequest(bytes_source=image_bytes)
            )
            result = poller.result()
            # Return extracted layout content (includes text, tables, and structure)
            return result.content.strip() if result and result.content else ""
        except Exception as e:
            raise RuntimeError(f"[Azure Fallback Error: {str(e)}]")

    # ==========================================================
    # CORE MULTIMODAL OCR (Gemini 1.5 Flash -> Azure Fallback)
    # ==========================================================
    def _run_multimodal_ocr(self, image_bytes: bytes) -> str:
        """
        Attempts OCR & visual analysis via Gemini-1.5-Flash.
        Falls back to Azure Document Intelligence if API limit (429) is hit.
        """
        if not image_bytes:
            return ""

        gemini = get_gemini_model()
        pil_img = Image.open(io.BytesIO(image_bytes))

        try:
            response = gemini.generate_content([self.vision_prompt, pil_img])
            if response and response.text:
                return response.text.strip()
            return ""
        except (google_exceptions.ResourceExhausted, google_exceptions.TooManyRequests) as rate_limit_err:
            # Explicit API limit / 429 quota exception: Fall back to Azure
            return self._run_azure_fallback(image_bytes)
        except Exception as e:
            # Check if the error message indicates rate limit/quota error
            err_str = str(e).lower()
            if "429" in err_str or "quota" in err_str or "resourceexhausted" in err_str:
                return self._run_azure_fallback(image_bytes)
            # Otherwise return error or fallback
            return self._run_azure_fallback(image_bytes)

    # ==========================================================
    # PARSER BRANCH A: TXT / MD
    # ==========================================================
    def _parse_text(self, file_path: str) -> str:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read().strip()

    # ==========================================================
    # PARSER BRANCH B: DOCUMENT PARSER (PDF, DOCX, PPTX)
    # ==========================================================
    def _parse_pdf(self, file_path: str) -> str:
        import fitz  # PyMuPDF
        doc = fitz.open(file_path)
        sections = []

        for page_num in range(len(doc)):
            page = doc[page_num]
            native_text = page.get_text().strip()

            # If page is scanned or image-heavy (less than 30 digital chars), render page to image
            if len(native_text) < 30:
                pix = page.get_pixmap(dpi=150)
                img_bytes = pix.tobytes("png")
                ocr_text = self._run_multimodal_ocr(img_bytes)
                if ocr_text:
                    sections.append(f"--- Page {page_num + 1} (Visual & Diagram OCR) ---\n{ocr_text}")
            else:
                sections.append(f"--- Page {page_num + 1} (Digital Text) ---\n{native_text}")

                # Process any diagrams/images embedded in the digital page
                for img_idx, img_info in enumerate(page.get_images(full=True)):
                    xref = img_info[0]
                    base_image = doc.extract_image(xref)
                    img_bytes = base_image["image"]
                    visual_ocr = self._run_multimodal_ocr(img_bytes)
                    if visual_ocr:
                        sections.append(f"[Page {page_num + 1} Diagram #{img_idx + 1}]:\n{visual_ocr}")

        doc.close()
        return "\n\n".join(sections)

    def _parse_docx(self, file_path: str) -> str:
        import docx
        doc = docx.Document(file_path)
        sections = []

        # 1. Native text: paragraphs and tables
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    paragraphs.append(row_text)
        if paragraphs:
            sections.append("\n".join(paragraphs))

        # 2. Embedded images and diagrams
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                img_bytes = rel.target_part.blob
                diagram_ocr = self._run_multimodal_ocr(img_bytes)
                if diagram_ocr:
                    sections.append(f"[Embedded Visual / Diagram]:\n{diagram_ocr}")

        return "\n\n".join(sections)

    def _parse_pptx(self, file_path: str) -> str:
        from pptx import Presentation
        prs = Presentation(file_path)
        sections = []

        for idx, slide in enumerate(prs.slides):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = "".join(run.text for run in paragraph.runs).strip()
                        if text:
                            slide_texts.append(text)
                # Embedded pictures, diagrams, and figures
                elif shape.shape_type == 13:  # Picture shape
                    img_bytes = shape.image.blob
                    visual_ocr = self._run_multimodal_ocr(img_bytes)
                    if visual_ocr:
                        slide_texts.append(f"[Slide Visual / Diagram]:\n{visual_ocr}")

            if slide_texts:
                sections.append(f"--- Slide {idx + 1} ---\n" + "\n".join(slide_texts))

        return "\n\n".join(sections)

    # ==========================================================
    # PARSER BRANCH C: STANDALONE IMAGES
    # ==========================================================
    def _parse_image(self, file_path: str) -> str:
        with open(file_path, "rb") as f:
            img_bytes = f.read()
        return self._run_multimodal_ocr(img_bytes)

    # ==========================================================
    # ENTRY POINT: MIME/EXTENSION DISPATCHER
    # ==========================================================
    def extract_text(self, file_or_path) -> str:
        """
        Dispatches to Text Parser, Document Parser, or Multimodal Image Parser
        and returns the unified text.
        """
        file_path = self._resolve_file_path(file_or_path)
        ext = os.path.splitext(file_path)[1].lower()
        mime_type, _ = mimetypes.guess_type(file_path)

        # Branch A: Plain text / Markdown
        if ext in ['.txt', '.md'] or mime_type in ['text/plain', 'text/markdown']:
            return self._parse_text(file_path)

        # Branch B: Document Parser
        if ext == '.pdf' or mime_type == 'application/pdf':
            return self._parse_pdf(file_path)
        if ext == '.docx' or mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
            return self._parse_docx(file_path)
        if ext == '.pptx' or mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
            return self._parse_pptx(file_path)

        # Branch C: Image Parser
        image_extensions = ['.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.webp']
        if ext in image_extensions or (mime_type and mime_type.startswith('image/')):
            return self._parse_image(file_path)

        # Fallback: Attempt image read; if invalid, attempt text read
        try:
            with open(file_path, "rb") as f:
                header = f.read(16)
            # If valid image header
            Image.open(file_path).verify()
            return self._parse_image(file_path)
        except Exception:
            try:
                return self._parse_text(file_path)
            except Exception:
                return ""